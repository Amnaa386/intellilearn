from datetime import datetime
from pathlib import Path
from typing import Dict, Any
from uuid import uuid4
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches
from pptx.util import Pt

from app.core.database import get_database
from app.models.presentations import PresentationCreate, PresentationResponse, PresentationTheme


class PresentationsService:
    THEME_COLORS = {
        PresentationTheme.CLASSIC: {
            "bg": RGBColor(255, 255, 255),
            "bg_alt": RGBColor(241, 245, 249),
            "accent": RGBColor(37, 99, 235),
            "text": RGBColor(15, 23, 42),
            "muted": RGBColor(71, 85, 105),
            "card": RGBColor(255, 255, 255),
            "card_border": RGBColor(203, 213, 225),
            "brand": "INTELLILEARN CLASSIC",
        },
        PresentationTheme.MODERN: {
            "bg": RGBColor(15, 23, 42),
            "bg_alt": RGBColor(30, 41, 59),
            "accent": RGBColor(56, 189, 248),
            "text": RGBColor(241, 245, 249),
            "muted": RGBColor(148, 163, 184),
            "card": RGBColor(30, 41, 59),
            "card_border": RGBColor(71, 85, 105),
            "brand": "INTELLILEARN MODERN AI",
        },
        PresentationTheme.PREMIUM: {
            "bg": RGBColor(23, 6, 53),
            "bg_alt": RGBColor(49, 24, 94),
            "accent": RGBColor(236, 72, 153),
            "text": RGBColor(250, 245, 255),
            "muted": RGBColor(196, 181, 253),
            "card": RGBColor(59, 29, 114),
            "card_border": RGBColor(167, 139, 250),
            "brand": "INTELLILEARN PREMIUM AI",
        },
    }

    @staticmethod
    def _clean_markdown_text(value: str) -> str:
        text = str(value or "").strip()
        if not text:
            return ""
        text = re.sub(r"`([^`]+)`", r"\1", text)
        text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
        text = re.sub(r"\*(.*?)\*", r"\1", text)
        text = re.sub(r"__(.*?)__", r"\1", text)
        text = re.sub(r"_(.*?)_", r"\1", text)
        text = re.sub(r"^\s*#{1,6}\s*", "", text)
        text = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1", text)
        text = re.sub(r"\s+", " ", text)
        return text.strip(" -:\t")

    @staticmethod
    def _normalize_section_title(title: str) -> str:
        cleaned = PresentationsService._clean_markdown_text(title)
        cleaned = re.sub(r"^(title|subtitle|slide)\s*:\s*", "", cleaned, flags=re.IGNORECASE)
        return cleaned[:100] or "Overview"

    @staticmethod
    def _require_db():
        db = get_database()
        if db is None:
            raise RuntimeError("Database connection is not available")
        return db

    @staticmethod
    def _sanitize_filename(value: str) -> str:
        cleaned = re.sub(r"[^a-zA-Z0-9\-_\s]", "", value or "").strip()
        cleaned = re.sub(r"\s+", "-", cleaned)
        return cleaned[:80] or "presentation"

    @staticmethod
    def _extract_sections(markdown: str):
        lines = str(markdown or "").replace("\r", "\n").split("\n")
        sections = []
        current_title = "Overview"
        current_points = []

        for raw in lines:
            line = raw.strip()
            if not line:
                continue
            if re.match(r"^#{1,3}\s+", line):
                if current_points:
                    sections.append((current_title, current_points))
                    current_points = []
                current_title = PresentationsService._normalize_section_title(re.sub(r"^#{1,3}\s+", "", line))
                continue
            if re.match(r"^[=\-]{3,}$", line):
                continue
            point = re.sub(r"^[-*]\s+", "", line)
            point = re.sub(r"^\d+\.\s+", "", point).strip()
            # Pull clean key/value when AI emits markdown labels.
            point = re.sub(r"^(title|subtitle)\s*:\s*", "", point, flags=re.IGNORECASE)
            point = PresentationsService._clean_markdown_text(point)
            if point:
                current_points.append(point)

        if current_points:
            sections.append((current_title, current_points))

        return sections[:12]

    def _build_ppt(self, title: str, topic: str, content: str, destination: Path, theme: PresentationTheme) -> int:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        palette = self.THEME_COLORS.get(theme, self.THEME_COLORS[PresentationTheme.MODERN])

        def paint_background(slide):
            bg = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0),
                Inches(0),
                prs.slide_width,
                prs.slide_height,
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = palette["bg"]
            bg.line.fill.background()

            top_band = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0),
                Inches(0),
                prs.slide_width,
                Inches(1.2),
            )
            top_band.fill.solid()
            top_band.fill.fore_color.rgb = palette["bg_alt"]
            top_band.line.fill.background()

            accent_bar = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.7),
                Inches(0.75),
                Inches(3.2),
                Inches(0.08),
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = palette["accent"]
            accent_bar.line.fill.background()

        def add_textbox(slide, left, top, width, height, text, size=20, bold=False, color=None, align=PP_ALIGN.LEFT):
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = align
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color or palette["text"]
            return box

        # Title slide
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        paint_background(slide)

        add_textbox(
            slide,
            Inches(0.8),
            Inches(0.25),
            Inches(8),
            Inches(0.45),
            palette["brand"],
            size=13,
            bold=True,
            color=palette["muted"],
        )

        add_textbox(
            slide,
            Inches(0.9),
            Inches(2.0),
            Inches(11.8),
            Inches(1.6),
            self._normalize_section_title(title),
            size=46,
            bold=True,
            color=palette["text"],
        )

        add_textbox(
            slide,
            Inches(0.95),
            Inches(4.2),
            Inches(11),
            Inches(0.6),
            f"Topic: {self._clean_markdown_text(topic)}",
            size=21,
            color=palette["muted"],
        )

        add_textbox(
            slide,
            Inches(0.95),
            Inches(4.85),
            Inches(8),
            Inches(0.45),
            "Generated by IntelliLearn AI",
            size=16,
            color=palette["accent"],
        )

        # Content slides
        sections = self._extract_sections(content)
        if not sections:
            sections = [("Key Points", [topic, "Auto-generated summary", "Review these concepts before practice"])]

        bullet_layout = prs.slide_layouts[6]
        for section_title, points in sections:
            section_slide = prs.slides.add_slide(bullet_layout)
            paint_background(section_slide)

            add_textbox(
                section_slide,
                Inches(0.8),
                Inches(0.24),
                Inches(11),
                Inches(0.5),
                self._normalize_section_title(section_title),
                size=28,
                bold=True,
                color=palette["text"],
            )

            card = section_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.8),
                Inches(1.35),
                Inches(11.8),
                Inches(5.6),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = palette["card"]
            card.line.color.rgb = palette["card_border"]
            card.line.width = Pt(1.4)

            body = section_slide.shapes.add_textbox(
                Inches(1.2),
                Inches(1.75),
                Inches(11.0),
                Inches(4.9),
            ).text_frame
            body.clear()
            body.word_wrap = True
            trimmed_points = points[:8]
            for idx, p in enumerate(trimmed_points):
                para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                para.text = f"• {self._clean_markdown_text(p[:220])}"
                para.level = 0
                para.font.size = Pt(22)
                para.font.color.rgb = palette["text"]
                para.space_after = Pt(11)

        destination.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(destination))
        return len(prs.slides)

    async def create_presentation(self, user_id: str, payload: PresentationCreate, public_base_url: str) -> PresentationResponse:
        presentation_id = str(uuid4())
        safe_name = self._sanitize_filename(payload.title)
        file_name = f"{safe_name}-{presentation_id[:8]}.pptx"
        output_path = Path("uploads") / "presentations" / file_name

        slide_count = self._build_ppt(payload.title, payload.topic, payload.content, output_path, payload.theme)
        file_url = f"{public_base_url}/uploads/presentations/{file_name}"

        doc = {
            "id": presentation_id,
            "userId": user_id,
            "title": payload.title,
            "topic": payload.topic,
            "fileName": file_name,
            "fileUrl": file_url,
            "slideCount": slide_count,
            "createdAt": datetime.utcnow(),
            "tags": ["AI-Generated", "Presentation", payload.topic],
            "theme": payload.theme.value,
        }

        db = self._require_db()
        db.collection("presentations").document(presentation_id).set(doc)
        return PresentationResponse(**doc)

    async def list_presentations(self, user_id: str, page: int = 1, limit: int = 50) -> Dict[str, Any]:
        db = self._require_db()
        skip = (page - 1) * limit
        docs = db.collection("presentations").where("userId", "==", user_id).stream()
        rows = []
        for doc in docs:
            row = doc.to_dict() or {}
            row["id"] = doc.id
            if not row.get("theme"):
                row["theme"] = PresentationTheme.MODERN.value
            rows.append(row)
        rows.sort(key=lambda x: x.get("createdAt") or datetime.min, reverse=True)
        total = len(rows)
        records = rows[skip:skip + limit]
        return {
            "presentations": [PresentationResponse(**item) for item in records],
            "total": total,
            "page": page,
            "limit": limit,
            "totalPages": (total + limit - 1) // limit,
        }

    async def delete_presentation(self, user_id: str, presentation_id: str) -> bool:
        db = self._require_db()
        ref = db.collection("presentations").document(presentation_id)
        snap = ref.get()
        if not snap.exists:
            return False
        row = snap.to_dict() or {}
        if row.get("userId") != user_id:
            return False
        file_name = row.get("fileName")
        if file_name:
            file_path = Path("uploads") / "presentations" / file_name
            if file_path.exists():
                file_path.unlink()
        ref.delete()
        return True


presentations_service = PresentationsService()
