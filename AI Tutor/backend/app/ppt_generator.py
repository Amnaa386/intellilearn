"""
PowerPoint Generation Service
Creates professional PowerPoint presentations using AI-generated content
"""

import os
import json
import tempfile
import logging
from typing import Dict, List, Optional, Any
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from .ai_services import AIService
from .config import settings

logger = logging.getLogger(__name__)

class PPTGenerator:
    """Service for generating PowerPoint presentations with AI content"""
    
    def __init__(self):
        self.ai_service = AIService()
        self.temp_dir = tempfile.gettempdir()
        logger.info("PPT Generator initialized")
    
    async def generate_presentation(self, topic: str, num_slides: int = 8) -> str:
        """
        Generate a complete PowerPoint presentation
        
        Args:
            topic: The main topic for the presentation
            num_slides: Number of slides to generate (default: 8)
            
        Returns:
            str: Path to the generated PPT file
        """
        logger.info(f"Generating presentation on topic: {topic}")
        
        try:
            # Generate structured content using AI
            slide_content = await self._generate_slide_content(topic, num_slides)
            logger.info(f"Generated content for {len(slide_content)} slides")
            
            # Create PowerPoint presentation
            ppt_path = await self._create_powerpoint(slide_content, topic)
            logger.info(f"PowerPoint created at: {ppt_path}")
            
            return ppt_path
            
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise Exception(f"Failed to generate presentation: {str(e)}")
    
    async def _generate_slide_content(self, topic: str, num_slides: int) -> List[Dict[str, Any]]:
        """
        Generate structured slide content using AI
        
        Args:
            topic: The main topic
            num_slides: Number of slides to generate
            
        Returns:
            List of slide content dictionaries
        """
        logger.info("Generating AI content for slides")
        
        # Create prompt for structured content generation
        prompt = f"""
        Generate a comprehensive PowerPoint presentation on "{topic}" with exactly {num_slides} slides.
        
        Format the response as a JSON array with the following structure:
        [
            {{
                "slide_number": 1,
                "title": "Slide Title",
                "content": [
                    "First bullet point",
                    "Second bullet point", 
                    "Third bullet point",
                    "Fourth bullet point",
                    "Fifth bullet point"
                ],
                "speaker_notes": "Detailed speaker notes for this slide",
                "slide_type": "title" or "content"
            }}
        ]
        
        Guidelines:
        1. Slide 1 should be a title slide with the main topic and subtitle
        2. Include an introduction slide (slide 2)
        3. Include 3-5 content slides covering key aspects
        4. Include a conclusion/summary slide
        5. Include a Q&A or thank you slide (last slide)
        6. Each content slide should have 3-5 bullet points
        7. Speaker notes should be detailed and helpful
        8. Content should be educational and professional
        
        Topic: {topic}
        Number of slides: {num_slides}
        
        Provide only the JSON response, no additional text.
        """
        
        try:
            # Use Gemini for structured content generation
            response = await self.ai_service.gemini_generate(prompt)
            
            if not response:
                logger.warning("Gemini failed, using fallback content")
                return self._get_fallback_content(topic, num_slides)
            
            # Parse JSON response
            try:
                # Clean up the response to extract JSON
                json_start = response.find('[')
                json_end = response.rfind(']') + 1
                
                if json_start == -1 or json_end == 0:
                    raise ValueError("No JSON array found in response")
                
                json_content = response[json_start:json_end]
                slide_data = json.loads(json_content)
                
                # Validate structure
                if not isinstance(slide_data, list) or len(slide_data) != num_slides:
                    raise ValueError(f"Expected {num_slides} slides, got {len(slide_data) if isinstance(slide_data, list) else 'invalid'}")
                
                # Validate each slide
                for i, slide in enumerate(slide_data):
                    if not all(key in slide for key in ['slide_number', 'title', 'content', 'speaker_notes']):
                        logger.warning(f"Slide {i+1} missing required fields, using fallback")
                        return self._get_fallback_content(topic, num_slides)
                
                logger.info("Successfully parsed AI-generated slide content")
                return slide_data
                
            except json.JSONDecodeError as e:
                logger.error(f"Failed to parse JSON response: {e}")
                return self._get_fallback_content(topic, num_slides)
                
        except Exception as e:
            logger.error(f"Error generating slide content: {e}")
            return self._get_fallback_content(topic, num_slides)
    
    def _get_fallback_content(self, topic: str, num_slides: int) -> List[Dict[str, Any]]:
        """Generate fallback content when AI fails"""
        logger.info("Using fallback content generation")
        
        slides = []
        
        # Title slide
        slides.append({
            "slide_number": 1,
            "title": topic,
            "content": [
                "An Educational Presentation",
                f"Generated on {datetime.now().strftime('%B %d, %Y')}",
                "AI Learning Platform",
                "Professional Content"
            ],
            "speaker_notes": f"Welcome to this presentation on {topic}. This slide introduces the main topic and sets the context for our discussion.",
            "slide_type": "title"
        })
        
        # Introduction slide
        slides.append({
            "slide_number": 2,
            "title": "Introduction",
            "content": [
                "Overview of the topic",
                "Key learning objectives",
                "Importance and relevance",
                "What we'll cover today"
            ],
            "speaker_notes": "This introduction provides an overview of what we'll be discussing and the key takeaways you can expect from this presentation.",
            "slide_type": "content"
        })
        
        # Content slides
        for i in range(3, num_slides - 1):
            slides.append({
                "slide_number": i,
                "title": f"Key Point {i-2}",
                "content": [
                    f"Important aspect of {topic}",
                    "Detailed explanation and examples",
                    "Practical applications",
                    "Key insights and learning",
                    "Related concepts"
                ],
                "speaker_notes": f"This slide covers key point {i-2} about {topic}. We'll explore the main concepts and their practical applications.",
                "slide_type": "content"
            })
        
        # Conclusion slide
        slides.append({
            "slide_number": num_slides,
            "title": "Summary & Q&A",
            "content": [
                "Key takeaways",
                "Summary of main points",
                "Further reading resources",
                "Questions and discussion"
            ],
            "speaker_notes": "Thank you for your attention! Let's summarize what we've learned and open the floor for any questions.",
            "slide_type": "content"
        })
        
        return slides
    
    async def _create_powerpoint(self, slide_content: List[Dict[str, Any]], topic: str) -> str:
        """
        Create PowerPoint presentation from structured content
        
        Args:
            slide_content: List of slide dictionaries
            topic: Main topic for filename
            
        Returns:
            str: Path to generated PPT file
        """
        logger.info("Creating PowerPoint presentation")
        
        # Create presentation object
        prs = Presentation()
        
        # Process each slide
        for slide_data in slide_content:
            await self._add_slide(prs, slide_data)
        
        # Generate unique filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = f"{safe_topic}_{timestamp}.pptx"
        
        # Create temporary file
        temp_path = os.path.join(self.temp_dir, filename)
        
        # Save presentation
        prs.save(temp_path)
        logger.info(f"Presentation saved to: {temp_path}")
        
        return temp_path
    
    async def _add_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add a single slide to the presentation"""
        
        slide_type = slide_data.get('slide_type', 'content')
        
        if slide_type == 'title':
            slide_layout = prs.slide_layouts[0]  # Title slide layout
        else:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = slide_data['title']
        
        # Configure title formatting
        title_font = title.text_frame.paragraphs[0].font
        title_font.name = 'Arial'
        title_font.size = Pt(36)
        title_font.bold = True
        title_font.color.rgb = RGBColor(0, 32, 96)  # Dark blue
        
        # Add content
        if slide_type != 'title' and 'content' in slide_data:
            content_shape = slide.placeholders[1]  # Content placeholder
            content_text_frame = content_shape.text_frame
            content_text_frame.clear()  # Clear existing text
            
            # Add bullet points
            for bullet_point in slide_data['content']:
                p = content_text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0  # Bullet level
                p.font.name = 'Arial'
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        # Add speaker notes
        if 'speaker_notes' in slide_data:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = slide_data['speaker_notes']
            
            # Format notes text
            for paragraph in notes_text_frame.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(12)
    
    def cleanup_file(self, file_path: str):
        """Clean up temporary file"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                logger.info(f"Cleaned up temporary file: {file_path}")
        except Exception as e:
            logger.error(f"Error cleaning up file {file_path}: {e}")
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get file information for download response"""
        try:
            stat = os.stat(file_path)
            return {
                "filename": os.path.basename(file_path),
                "size": stat.st_size,
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "path": file_path
            }
        except Exception as e:
            logger.error(f"Error getting file info: {e}")
            return {}
